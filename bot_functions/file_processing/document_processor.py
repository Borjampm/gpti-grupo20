import os
import tempfile
import pandas as pd
from docx import Document
from pptx import Presentation
import subprocess
import platform

async def convert_docx_to_pdf(docx_path: str, chat_id: int) -> str:
    """Convert DOCX file to PDF with formatting preservation"""
    try:
        temp_dir = tempfile.gettempdir()
        output_path = os.path.join(temp_dir, f"docx_to_pdf_{chat_id}.pdf")

        # Try docx2pdf first for best formatting preservation
        try:
            from docx2pdf import convert

            # Convert using docx2pdf
            convert(docx_path, output_path)

            if os.path.exists(output_path):
                print(f"Successfully converted using docx2pdf: {docx_path}")
                return output_path

        except ImportError:
            print("docx2pdf not available, falling back to alternative methods")
        except Exception as e:
            print(f"docx2pdf conversion failed: {e}, falling back to alternative methods")

        # Fallback to LibreOffice for Linux/Mac
        try:
            # Try both LibreOffice command names
            libreoffice_commands = ['libreoffice', 'soffice']
            for cmd in libreoffice_commands:
                try:
                    subprocess.run([
                        cmd, '--headless', '--convert-to', 'pdf',
                        '--outdir', temp_dir, docx_path
                    ], check=True, capture_output=True)

                    # LibreOffice creates file with same name but .pdf extension
                    base_name = os.path.splitext(os.path.basename(docx_path))[0]
                    libreoffice_output = os.path.join(temp_dir, f"{base_name}.pdf")

                    if os.path.exists(libreoffice_output):
                        os.rename(libreoffice_output, output_path)
                        print(f"Successfully converted using {cmd}: {docx_path}")
                        return output_path
                except FileNotFoundError:
                    continue  # Try next command
                except subprocess.CalledProcessError:
                    continue  # Try next command

        except Exception as e:
            print(f"LibreOffice conversion failed: {e}")

        # Final fallback: create simple PDF from text
        print("Using text extraction fallback method")
        return await _create_simple_pdf_from_docx(docx_path, chat_id)

    except Exception as e:
        print(f"Error converting DOCX to PDF: {e}")
        return None

async def _create_simple_pdf_from_docx(docx_path: str, chat_id: int) -> str:
    """Fallback method to create PDF from DOCX using text extraction"""
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter

        temp_dir = tempfile.gettempdir()
        output_path = os.path.join(temp_dir, f"docx_to_pdf_simple_{chat_id}.pdf")

        # Extract text from DOCX
        doc = Document(docx_path)
        text_content = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)

        # Create PDF
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        y_position = height - 50

        for line in text_content:
            if y_position < 50:  # Start new page
                c.showPage()
                y_position = height - 50

            # Wrap long lines
            words = line.split()
            current_line = ""
            for word in words:
                test_line = f"{current_line} {word}".strip()
                if len(test_line) * 6 < width - 100:  # Rough character width estimation
                    current_line = test_line
                else:
                    if current_line:
                        c.drawString(50, y_position, current_line)
                        y_position -= 15
                    current_line = word

            if current_line:
                c.drawString(50, y_position, current_line)
                y_position -= 20

        c.save()
        print(f"Created simple text-based PDF: {output_path}")
        return output_path
    except Exception as e:
        print(f"Error creating simple PDF from DOCX: {e}")
        return None

async def convert_pdf_to_docx(pdf_path: str, chat_id: int) -> str:
    """Convert PDF file to DOCX with formatting preservation"""
    try:
        temp_dir = tempfile.gettempdir()
        output_path = os.path.join(temp_dir, f"pdf_to_docx_{chat_id}.docx")

        # Try pdf2docx first for best formatting preservation
        try:
            from pdf2docx import Converter

            # Create converter instance
            cv = Converter(pdf_path)

            # Convert PDF to DOCX with formatting preservation
            cv.convert(output_path, start=0, end=None)
            cv.close()

            if os.path.exists(output_path):
                print(f"Successfully converted using pdf2docx: {pdf_path}")
                return output_path

        except ImportError:
            print("pdf2docx not available, falling back to text extraction method")
        except Exception as e:
            print(f"pdf2docx conversion failed: {e}, falling back to text extraction method")

        # Fallback to PyPDF2 text extraction method
        print("Using PyPDF2 text extraction fallback method")
        return await _create_simple_docx_from_pdf(pdf_path, chat_id)

    except Exception as e:
        print(f"Error converting PDF to DOCX: {e}")
        return None

async def _create_simple_docx_from_pdf(pdf_path: str, chat_id: int) -> str:
    """Fallback method to create DOCX from PDF using text extraction"""
    try:
        from PyPDF2 import PdfReader
        from docx import Document
        from docx.shared import Inches
        from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

        temp_dir = tempfile.gettempdir()
        output_path = os.path.join(temp_dir, f"pdf_to_docx_simple_{chat_id}.docx")

        # Extract text from PDF
        reader = PdfReader(pdf_path)

        # Create DOCX document with better formatting
        doc = Document()

        # Add title
        title = doc.add_heading('Documento convertido desde PDF', 0)
        title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        # Add some spacing
        doc.add_paragraph()

        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text.strip():
                # Add page header
                if page_num > 1:
                    doc.add_page_break()

                page_header = doc.add_heading(f'Página {page_num}', level=2)
                page_header.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT

                # Process text into paragraphs
                paragraphs = text.split('\n\n')
                for para_text in paragraphs:
                    if para_text.strip():
                        # Clean up the text
                        cleaned_text = ' '.join(para_text.split())
                        if len(cleaned_text) > 10:  # Only add substantial paragraphs
                            para = doc.add_paragraph(cleaned_text)
                            para.style = 'Normal'

                # Add some spacing between pages
                doc.add_paragraph()
            else:
                # Handle pages with no extractable text
                if page_num > 1:
                    doc.add_page_break()

                page_header = doc.add_heading(f'Página {page_num}', level=2)
                no_text_para = doc.add_paragraph('[Esta página no contiene texto extraíble o contiene solo imágenes]')
                no_text_para.italic = True

        # Add footer note
        doc.add_paragraph()
        footer_para = doc.add_paragraph()
        footer_run = footer_para.add_run('Nota: Este documento fue convertido automáticamente desde PDF. '
                                       'La formatación original puede no estar completamente preservada.')
        footer_run.italic = True
        footer_run.font.size = 12
        footer_para.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        doc.save(output_path)
        print(f"Created enhanced DOCX from PDF text extraction: {output_path}")
        return output_path

    except Exception as e:
        print(f"Error creating simple DOCX from PDF: {e}")
        return None

async def convert_csv_to_excel(csv_path: str, chat_id: int) -> str:
    """Convert CSV file to Excel"""
    try:
        temp_dir = tempfile.gettempdir()
        output_path = os.path.join(temp_dir, f"csv_to_excel_{chat_id}.xlsx")

        # Read CSV and convert to Excel
        df = pd.read_csv(csv_path)
        df.to_excel(output_path, index=False)

        return output_path
    except Exception as e:
        print(f"Error converting CSV to Excel: {e}")
        return None

async def convert_excel_to_csv(excel_path: str, chat_id: int) -> str:
    """Convert Excel file to CSV"""
    try:
        temp_dir = tempfile.gettempdir()
        output_path = os.path.join(temp_dir, f"excel_to_csv_{chat_id}.csv")

        # Read Excel and convert to CSV
        df = pd.read_excel(excel_path, sheet_name=0)  # Read first sheet
        df.to_csv(output_path, index=False)

        return output_path
    except Exception as e:
        print(f"Error converting Excel to CSV: {e}")
        return None

async def convert_pptx_to_pdf(pptx_path: str, chat_id: int) -> str:
    """Convert PowerPoint file to PDF with formatting preservation using LibreOffice CLI"""
    try:
        temp_dir = tempfile.gettempdir()
        output_path = os.path.join(temp_dir, f"pptx_to_pdf_{chat_id}.pdf")

        # Try LibreOffice CLI first (cross-platform, best formatting preservation)
        # Try both command names: 'libreoffice' and 'soffice'
        libreoffice_commands = ['libreoffice', 'soffice']
        libreoffice_success = False

        for cmd in libreoffice_commands:
            try:
                # Use LibreOffice headless mode for conversion
                result = subprocess.run([
                    cmd, '--headless', '--convert-to', 'pdf',
                    '--outdir', temp_dir, pptx_path
                ], check=True, capture_output=True, text=True, timeout=60)

                # LibreOffice creates file with same name but .pdf extension
                base_name = os.path.splitext(os.path.basename(pptx_path))[0]
                libreoffice_output = os.path.join(temp_dir, f"{base_name}.pdf")

                if os.path.exists(libreoffice_output):
                    # Rename to our desired output path
                    if libreoffice_output != output_path:
                        os.rename(libreoffice_output, output_path)
                    print(f"Successfully converted using {cmd} CLI: {pptx_path}")
                    libreoffice_success = True
                    break
                else:
                    print(f"{cmd} conversion completed but output file not found")

            except subprocess.TimeoutExpired:
                print(f"{cmd} conversion timed out (60s limit)")
            except subprocess.CalledProcessError as e:
                print(f"{cmd} conversion failed with exit code {e.returncode}: {e.stderr}")
            except FileNotFoundError:
                print(f"{cmd} not found in PATH")
                continue  # Try next command
            except Exception as e:
                print(f"{cmd} conversion failed: {e}")

        if libreoffice_success:
            return output_path

        # Fallback for Windows: try comtypes (Windows only)
        if platform.system() == "Windows":
            try:
                import comtypes.client
                print("Trying Windows PowerPoint COM interface...")

                powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                powerpoint.Visible = 1

                presentation = powerpoint.Presentations.Open(os.path.abspath(pptx_path))
                presentation.SaveAs(os.path.abspath(output_path), 32)  # 32 = PDF format
                presentation.Close()
                powerpoint.Quit()

                if os.path.exists(output_path):
                    print(f"Successfully converted using Windows PowerPoint COM: {pptx_path}")
                    return output_path

            except ImportError:
                print("comtypes not available for Windows PowerPoint conversion")
            except Exception as e:
                print(f"Windows PowerPoint COM conversion failed: {e}")

        # Final fallback: create simple PDF from presentation text
        print("Using text extraction fallback method for PowerPoint")
        return await _create_simple_pdf_from_pptx(pptx_path, chat_id)

    except Exception as e:
        print(f"Error converting PPTX to PDF: {e}")
        return None

async def _create_simple_pdf_from_pptx(pptx_path: str, chat_id: int) -> str:
    """Fallback method to create PDF from PPTX using text extraction"""
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.units import inch

        temp_dir = tempfile.gettempdir()
        output_path = os.path.join(temp_dir, f"pptx_to_pdf_simple_{chat_id}.pdf")

        # Extract text from PPTX
        prs = Presentation(pptx_path)
        slides_content = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            if slide_text:
                slides_content.append({
                    'number': slide_num,
                    'content': '\n'.join(slide_text)
                })
            else:
                slides_content.append({
                    'number': slide_num,
                    'content': f"[Slide {slide_num} - No extractable text content]"
                })

        if not slides_content:
            slides_content.append({
                'number': 1,
                'content': "[No slides or content found in presentation]"
            })

        # Create PDF with better formatting
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        margin = 0.75 * inch

        for i, slide_data in enumerate(slides_content):
            if i > 0:
                c.showPage()

            # Title for each slide
            y_position = height - margin
            c.setFont("Helvetica-Bold", 16)
            title = f"Slide {slide_data['number']}"
            c.drawString(margin, y_position, title)
            y_position -= 30

            # Draw a line under the title
            c.setLineWidth(1)
            c.line(margin, y_position, width - margin, y_position)
            y_position -= 20

            # Content
            c.setFont("Helvetica", 12)
            lines = slide_data['content'].split('\n')

            for line in lines:
                if y_position < margin + 50:  # Near bottom of page
                    c.showPage()
                    y_position = height - margin

                # Handle long lines by wrapping
                words = line.split()
                current_line = ""
                max_width = width - (2 * margin)

                for word in words:
                    test_line = f"{current_line} {word}".strip()
                    # Rough estimation: 6 points per character
                    if len(test_line) * 6 < max_width:
                        current_line = test_line
                    else:
                        if current_line:
                            c.drawString(margin, y_position, current_line)
                            y_position -= 15
                        current_line = word

                if current_line:
                    c.drawString(margin, y_position, current_line)
                    y_position -= 15

                # Extra space between text blocks
                y_position -= 5

        c.save()
        print(f"Created simple text-based PDF from PowerPoint: {output_path}")
        return output_path

    except Exception as e:
        print(f"Error creating simple PDF from PPTX: {e}")
        return None
